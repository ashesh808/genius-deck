from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import io
from modules.Image_caption import Image_Caption
import os

class Powerpoint(Image_Caption):
    def __init__(self, Upload_path, id, Flag=bool):
        self.pptx_path = os.path.join(Upload_path, id + '.pptx')
        self.Flag = Flag
        self.save = os.path.join(Upload_path, id + '.txt')
        Image_Caption.__init__(self)

    def PptToText(self):
        whole_text = ''
        prs = Presentation(self.pptx_path)
       #f = open(self.save, 'a') #change to final location
        for slide_number, slide in enumerate(prs.slides):
            #print(f"Slide {slide_number + 1}:")
            x = str(slide_number)
            #f.write("Slide " + x + "\n")
            whole_text += "Slide " + x + "\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    #f.write(shape.text)
                    whole_text += shape.text
                    #print(shape.text)
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE and self.Flag != 'True':
                    txt = 'An image depicting '
                    Image = shape.image
                    Image_bytes = Image.blob
                    self.Img = io.BytesIO(Image_bytes)
                    Powerpoint.Caption(self)
                    text = Powerpoint.ListToString(self)
                    text.strip('['']')
                    text = txt + text + '\n'
                    #f.write(text)
                    whole_text += text
        return whole_text
        #print(whole_text)
        #f.close()

if __name__ == '__main__':
    pptx = input('Input path ')
    test = Powerpoint(pptx, False)
    test.PptToText()
