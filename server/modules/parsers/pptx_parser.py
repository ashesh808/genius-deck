from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import io
import os

from modules.parsers.Image_caption import Image_Caption
from modules.parsers.IParser import IParser

class PptxParser(IParser, Image_Caption):
    def __init__(self, Upload_path, id, Flag=True):
        self.pptx_path = os.path.join(Upload_path, id + '.pptx')
        self.Flag = Flag
        self.save = os.path.join(Upload_path, id + '.txt')
        Image_Caption.__init__(self)

    def parse(self):
        whole_text = ''
        prs = Presentation(self.pptx_path)
        for slide_number, slide in enumerate(prs.slides):
            x = str(slide_number)
            whole_text += "Slide " + x + "\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    whole_text += shape.text
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE and self.Flag != 'True':
                    txt = 'An image depicting '
                    Image = shape.image
                    Image_bytes = Image.blob
                    self.Img = io.BytesIO(Image_bytes)
                    PptxParser.Caption(self)
                    text = PptxParser.ListToString(self)
                    text.strip('['']')
                    text = txt + text + '\n'
                    whole_text += text
        return whole_text

